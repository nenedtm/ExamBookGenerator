"""PPTX parser for ExamBookGenerator.

Extracts text (slide titles, body, notes) and embedded images from
PowerPoint presentations using python-pptx.  Images are delegated to
``pipeline.image_extractor.save_extracted_image`` for validation,
deduplication, and disk storage.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE, MSO_SHAPE_TYPE

from core.models import Document, ExtractedImage, FileType
from pipeline.image_extractor import save_extracted_image
from utils.logger import get_logger

logger = get_logger(__name__)


def _slide_title(slide: object) -> str | None:
    """Return the title text of a slide, or *None*."""
    if not hasattr(slide, "shapes"):
        return None
    for shape in slide.shapes:  # type: ignore[union-attr]
        if not getattr(shape, "has_text_frame", False):
            continue
        # Placeholder title shapes have type TITLE or CENTER_TITLE
        if getattr(shape, "placeholder_format", None) is not None:
            ptype = shape.placeholder_format.type  # type: ignore[union-attr]
            if ptype in (PP_PLACEHOLDER_TYPE.TITLE, PP_PLACEHOLDER_TYPE.CENTER_TITLE):
                text = shape.text_frame.text.strip()  # type: ignore[union-attr]
                if text:
                    return text
    # Fallback: first text shape
    for shape in slide.shapes:  # type: ignore[union-attr]
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()  # type: ignore[union-attr]
            if text:
                return text
    return None


def _slide_body_text(slide: object) -> str:
    """Return all text content of a slide (body text, excluding title)."""
    parts: list[str] = []
    if not hasattr(slide, "shapes"):
        return ""
    for shape in slide.shapes:  # type: ignore[union-attr]
        if not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "placeholder_format", None) is not None:
            ptype = shape.placeholder_format.type  # type: ignore[union-attr]
            if ptype in (PP_PLACEHOLDER_TYPE.TITLE, PP_PLACEHOLDER_TYPE.CENTER_TITLE):
                continue
        text = shape.text_frame.text.strip()  # type: ignore[union-attr]
        if text:
            parts.append(text)
    return "\n".join(parts)


def _slide_notes(slide: object) -> str | None:
    """Return the speaker notes of a slide, or *None*."""
    if not getattr(slide, "has_notes_slide", False):
        return None
    try:
        notes_frame = slide.notes_slide.notes_text_frame  # type: ignore[union-attr]
        text = notes_frame.text.strip()
        return text if text else None
    except Exception:
        return None


def _extract_images(
    prs: Presentation,
    source_document_id: str,
) -> list[ExtractedImage]:
    """Iterate every slide, find PICTURE shapes, and delegate extraction."""
    images: list[ExtractedImage] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        if not hasattr(slide, "shapes"):
            continue
        for shape in slide.shapes:  # type: ignore[union-attr]
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:  # type: ignore[union-attr]
                    continue
            except Exception:
                continue

            try:
                raw_bytes = shape.image.blob  # type: ignore[union-attr]
            except Exception:
                logger.warning(
                    "Cannot read image from slide %d in doc %s — skipped",
                    slide_idx, source_document_id,
                )
                continue

            if not raw_bytes:
                continue

            # Build a caption from alt text / name
            caption: str | None = None
            if hasattr(shape, "name") and shape.name:  # type: ignore[union-attr]
                caption = shape.name  # type: ignore[union-attr]

            extracted = save_extracted_image(
                raw_bytes=raw_bytes,
                source_document_id=source_document_id,
                page_or_slide=slide_idx,
                caption=caption,
            )
            if extracted is not None:
                images.append(extracted)

    return images


def parse_pptx(path: Path | str) -> tuple[Document, list[ExtractedImage]]:
    """Parse a PPTX file, extracting text, notes, and embedded images.

    Parameters
    ----------
    path:
        Path to a ``.pptx`` file.

    Returns
    -------
    tuple[Document, list[ExtractedImage]]
        A ``Document`` with ``content``, ``metadata``, and ``images``
        populated, plus the list of ``ExtractedImage`` objects saved to
        disk.

    Raises
    ------
    FileNotFoundError
        If *path* does not exist.
    """
    pptx_path = Path(path)

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    doc = Document(source_path=str(pptx_path), file_type=FileType.PPTX)
    images: list[ExtractedImage] = []

    try:
        prs = Presentation(pptx_path)
    except Exception:
        logger.warning("Cannot open PPTX %s — corrupted or unreadable", pptx_path)
        return doc, images

    # ── Metadata ────────────────────────────────────────────────────────
    core = prs.core_properties
    meta_fields = {
        "title": core.title,
        "author": core.author,
        "subject": core.subject,
        "category": core.category,
    }
    for key, val in meta_fields.items():
        if val:
            doc.metadata[key] = val

    doc.metadata["slide_count"] = len(prs.slides)

    if core.title:
        doc.title = core.title

    # ── Slide iteration ─────────────────────────────────────────────────
    slide_texts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        body = _slide_body_text(slide)
        notes = _slide_notes(slide)

        parts: list[str] = []
        if title:
            parts.append(f"## Slide {slide_idx}: {title}")
        else:
            parts.append(f"## Slide {slide_idx}")

        if body:
            parts.append(body)
        if notes:
            parts.append(f"[Notes]: {notes}")

        slide_text = "\n".join(parts)
        if slide_text.strip():
            slide_texts.append(slide_text)

    # ── Images ──────────────────────────────────────────────────────────
    images = _extract_images(prs, doc.id)
    for img in images:
        doc.images.append(img.id)

    doc.content = "\n\n".join(slide_texts)

    if not doc.content.strip():
        logger.info("PPTX %s yielded no extractable text", pptx_path)
    else:
        logger.info(
            "Parsed PPTX %s — %d chars, %d image(s)",
            pptx_path, len(doc.content), len(images),
        )

    return doc, images
