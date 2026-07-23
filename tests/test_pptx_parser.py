"""Tests for parsers.pptx_parser."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from core.models import Document, FileType
from parsers.pptx_parser import parse_pptx


# ── Helpers ─────────────────────────────────────────────────────────────────

def _create_pptx(
    path: Path,
    slides_text: list[list[str]] | None = None,
    title: str = "",
    add_image: bool = False,
    add_notes: bool = False,
) -> Path:
    """Build a minimal PPTX for testing."""
    prs = Presentation()

    if slides_text is None:
        slides_text = [["Hello world"]]

    for slide_idx, texts in enumerate(slides_text):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

        # Title
        title_shape = slide.placeholders[0]
        title_shape.text = texts[0] if texts else f"Slide {slide_idx + 1}"

        # Body text
        if len(texts) > 1:
            body_shape = slide.placeholders[1]
            body_shape.text = "\n".join(texts[1:])

        if add_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = f"Notes for slide {slide_idx + 1}"

        if add_image:
            img = Image.new("RGB", (200, 150), color="blue")
            buf = BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            slide.shapes.add_picture(buf, Inches(1), Inches(1), Inches(2), Inches(1.5))

    if title:
        prs.core_properties.title = title

    prs.save(path)
    return path


def _create_pptx_empty(path: Path) -> Path:
    """Create a PPTX with no slides."""
    prs = Presentation()
    prs.save(path)
    return path


def _corrupt_pptx(path: Path) -> Path:
    path.write_bytes(b"PK\x03\x04corrupt zip data here")
    return path


# ── parse_pptx ──────────────────────────────────────────────────────────────

class TestParsePptx:
    def test_returns_document(self, tmp_path: Path) -> None:
        p = _create_pptx(tmp_path / "test.pptx")
        doc, images = parse_pptx(p)
        assert isinstance(doc, Document)
        assert isinstance(images, list)

    def test_file_type_is_pptx(self, tmp_path: Path) -> None:
        p = _create_pptx(tmp_path / "test.pptx")
        doc, _ = parse_pptx(p)
        assert doc.file_type == FileType.PPTX

    def test_source_path_matches(self, tmp_path: Path) -> None:
        p = _create_pptx(tmp_path / "test.pptx")
        doc, _ = parse_pptx(p)
        assert doc.source_path == str(p)

    def test_text_extracted(self, tmp_path: Path) -> None:
        p = _create_pptx(tmp_path / "test.pptx", slides_text=[["Title", "Body text here"]])
        doc, _ = parse_pptx(p)
        assert "Title" in doc.content
        assert "Body text here" in doc.content

    def test_multiple_slides(self, tmp_path: Path) -> None:
        p = _create_pptx(
            tmp_path / "test.pptx",
            slides_text=[["First slide", "First body"], ["Second slide", "Second body"]],
        )
        doc, _ = parse_pptx(p)
        assert "First slide" in doc.content
        assert "Second slide" in doc.content
        assert "First body" in doc.content
        assert "Second body" in doc.content

    def test_slide_numbering(self, tmp_path: Path) -> None:
        p = _create_pptx(tmp_path / "test.pptx", slides_text=[["A"], ["B"]])
        doc, _ = parse_pptx(p)
        assert "Slide 1" in doc.content
        assert "Slide 2" in doc.content

    def test_metadata_title(self, tmp_path: Path) -> None:
        p = _create_pptx(tmp_path / "test.pptx", title="My Presentation")
        doc, _ = parse_pptx(p)
        assert doc.title == "My Presentation"
        assert doc.metadata.get("title") == "My Presentation"

    def test_slide_count_metadata(self, tmp_path: Path) -> None:
        p = _create_pptx(
            tmp_path / "test.pptx",
            slides_text=[["A"], ["B"], ["C"]],
        )
        doc, _ = parse_pptx(p)
        assert doc.metadata.get("slide_count") == 3

    def test_notes_extracted(self, tmp_path: Path) -> None:
        p = _create_pptx(
            tmp_path / "test.pptx",
            slides_text=[["Title", "Body"]],
            add_notes=True,
        )
        doc, _ = parse_pptx(p)
        assert "[Notes]:" in doc.content
        assert "Notes for slide 1" in doc.content

    def test_images_extracted(self, tmp_path: Path) -> None:
        p = _create_pptx(
            tmp_path / "test.pptx",
            slides_text=[["Title"]],
            add_image=True,
        )
        doc, images = parse_pptx(p)
        assert len(images) >= 1
        assert len(doc.images) >= 1

    def test_image_file_on_disk(self, tmp_path: Path) -> None:
        p = _create_pptx(
            tmp_path / "test.pptx",
            slides_text=[["Title"]],
            add_image=True,
        )
        _, images = parse_pptx(p)
        assert Path(images[0].file_path).exists()

    def test_image_slide_number(self, tmp_path: Path) -> None:
        p = _create_pptx(
            tmp_path / "test.pptx",
            slides_text=[["Slide 1"]],
            add_image=True,
        )
        _, images = parse_pptx(p)
        # Image was added to slide 1 (only slide)
        assert images[0].page_or_slide == 1

    def test_no_images_pptx(self, tmp_path: Path) -> None:
        p = _create_pptx(
            tmp_path / "test.pptx",
            slides_text=[["Title"]],
            add_image=False,
        )
        doc, images = parse_pptx(p)
        assert images == []
        assert doc.images == []

    def test_empty_pptx(self, tmp_path: Path) -> None:
        p = _create_pptx_empty(tmp_path / "empty.pptx")
        doc, images = parse_pptx(p)
        assert doc.content.strip() == ""
        assert images == []
        assert doc.metadata.get("slide_count") == 0

    def test_empty_slide_text(self, tmp_path: Path) -> None:
        p = _create_pptx(tmp_path / "test.pptx", slides_text=[[""]])
        doc, _ = parse_pptx(p)
        # Should still have slide numbering even if text is empty
        assert "Slide 1" in doc.content

    def test_corrupt_pptx(self, tmp_path: Path) -> None:
        _corrupt_pptx(tmp_path / "bad.pptx")
        doc, images = parse_pptx(tmp_path / "bad.pptx")
        assert doc.content == ""
        assert images == []

    def test_nonexistent_raises(self) -> None:
        with pytest.raises(FileNotFoundError):
            parse_pptx(Path("/no/such/file.pptx"))
